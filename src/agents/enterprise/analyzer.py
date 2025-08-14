"""
Enterprise document analyzer for creative consultation.
One-time analysis of business documents to generate comprehensive creative briefs.
"""

import json
from typing import Dict, Any, Optional, Callable
from datetime import datetime
from pathlib import Path

from google import genai
from google.genai import types
from ...core.config import PROJECT_ID, CREDENTIALS
from ...core.state import RAGState


# Concise system instruction for LLM role
ENTERPRISE_SYSTEM_INSTRUCTION = """You are a strategic planner at a 4A creative agency.

Your role is to read client business materials and translate them into strategy briefs for the creative team. You explore branding and marketing strategies, extracting insights about positioning, audience needs, and messaging frameworks.

You are NOT on the creative team, so your analysis does not dictate specific creative treatments or execution techniques. Instead, you provide strategic context that enables the creative team to craft their own creative vision.

Focus on actionable strategic insights about business context, audience understanding, and brand positioning."""


# Detailed user prompt sent with documents
ENTERPRISE_ANALYSIS_PROMPT = """Analyze the provided documents and create a succinct creative brief:

1. **Business Context**: What is this company/brand about? What do they do? Summarize useful information from materials for strategizing. (1-5 sentences)

2. **Products & Services**: What are they offering? Key features and differentiators? How do products work or transform (foldable, raising components, lighting behavior, etc.)? Any features that should NOT be shown? Specific visualization accuracy requirements? (1-5 sentences)

3. **Brand Position & Promise**: For [target audience], [brand] is the [category] that [point of difference]. Promise: [core value proposition]. (1-3 sentences)

4. **Target Audience & Insight**: Who are they? What is the core human truth or tension they experience? What job are they hiring this product/service to do? (1-3 sentences)

5. **Message Strategy**: What are the key themes, tone, and decision drivers? What proof points support the claims? (1-3 sentences)

6. **Creative Opportunities**: What video content would be most impactful? What stories, visual styles, and approaches work best? (1-3 sentences)

7. **Constraints**: Any specific requirements or guidelines? (1-3 sentences, or "None specified")

Write in professional and natural language. Focus on actionable insights for video creators.

If documents contain visual identity elements, note them briefly."""


def analyze_enterprise_documents(
    state: RAGState,
    stream_callback: Optional[Callable] = None
) -> Dict[str, Any]:
    """
    Analyze enterprise documents to generate creative brief.

    Reads all documents from state["enterprise_resources"]["documents"],
    uses Gemini 2.5 Pro for comprehensive analysis,
    returns structured output for creative workflow.

    Args:
        state: RAGState with enterprise_resources
        stream_callback: Optional callback for progress events

    Returns:
        Analysis output dict to be saved in state["enterprise_agent_output"]
    """
    print("[Enterprise Analyzer] Starting document analysis...")
    start_time = datetime.now()

    enterprise_resources = state.get("enterprise_resources", {})
    documents = enterprise_resources.get("documents", [])
    images = enterprise_resources.get("images", [])

    if not documents and not images:
        print("[Enterprise Analyzer] No resources to analyze")
        return {
            "status": "no_resources",
            "message": "No enterprise resources uploaded yet",
            "timestamp": datetime.now().isoformat()
        }

    print(f"[Enterprise Analyzer] Analyzing {len(documents)} documents, {len(images)} images")

    # Emit start event
    if state:
        from ..base import emit_event
        emit_event(state, "enterprise_analysis_started", {
            "document_count": len(documents),
            "image_count": len(images),
            "timestamp": datetime.now().isoformat()
        }, agent_name="enterprise_analyzer")

    try:
        # Initialize Gemini client
        client = genai.Client(
            vertexai=True,
            project=PROJECT_ID,
            location="global",
            credentials=CREDENTIALS
        )

        # Check GCS download utility availability
        try:
            from ...storage.gcs_utils import download_to_temp
            has_gcs = True
        except ImportError:
            has_gcs = False
            print("[Enterprise Analyzer] Warning: GCS utils not available")

        # Build content parts for Gemini
        content_parts = []
        processed_files = []

        # Process documents
        for doc in documents:
            url = doc.get("url", "")
            filename = doc.get("filename", "unknown")
            extension = doc.get("extension", "")

            if not url:
                continue

            try:
                if has_gcs and (url.startswith("gs://") or url.startswith("http")):
                    # Download and process
                    with download_to_temp(url, suffix=extension) as temp_path:
                        with open(temp_path, 'rb') as f:
                            doc_bytes = f.read()

                        # Add to content based on type
                        if extension == ".pdf":
                            part = types.Part.from_bytes(
                                data=doc_bytes,
                                mime_type="application/pdf"
                            )
                            content_parts.append(part)
                            processed_files.append({"filename": filename, "type": "pdf"})
                            print(f"[Enterprise Analyzer] Loaded PDF: {filename}")

                        elif extension in [".pptx", ".ppt"]:
                            # Try to extract text from PowerPoint
                            try:
                                from pptx import Presentation
                                import tempfile

                                # Save to temp file for pptx library
                                with tempfile.NamedTemporaryFile(suffix=extension, delete=False) as tmp:
                                    tmp.write(doc_bytes)
                                    tmp_path = tmp.name

                                prs = Presentation(tmp_path)
                                text_content = []
                                for slide in prs.slides:
                                    for shape in slide.shapes:
                                        if hasattr(shape, "text"):
                                            text_content.append(shape.text)

                                import os
                                os.unlink(tmp_path)

                                # Add as text
                                content_parts.append("\n\n".join(text_content))
                                processed_files.append({"filename": filename, "type": "ppt"})
                                print(f"[Enterprise Analyzer] Extracted text from PPT: {filename}")

                            except ImportError:
                                print(f"[Enterprise Analyzer] python-pptx not available, skipping {filename}")

                        elif extension in [".docx", ".doc"]:
                            # Try to extract text from Word
                            try:
                                from docx import Document
                                import tempfile

                                with tempfile.NamedTemporaryFile(suffix=extension, delete=False) as tmp:
                                    tmp.write(doc_bytes)
                                    tmp_path = tmp.name

                                doc_obj = Document(tmp_path)
                                text_content = [para.text for para in doc_obj.paragraphs]

                                import os
                                os.unlink(tmp_path)

                                # Add as text
                                content_parts.append("\n\n".join(text_content))
                                processed_files.append({"filename": filename, "type": "docx"})
                                print(f"[Enterprise Analyzer] Extracted text from Word: {filename}")

                            except ImportError:
                                print(f"[Enterprise Analyzer] python-docx not available, skipping {filename}")

                else:
                    # Local file path (for testing)
                    if extension == ".pdf":
                        with open(url, 'rb') as f:
                            doc_bytes = f.read()
                        part = types.Part.from_bytes(data=doc_bytes, mime_type="application/pdf")
                        content_parts.append(part)
                        processed_files.append({"filename": filename, "type": "pdf"})

            except Exception as e:
                print(f"[Enterprise Analyzer] Error processing {filename}: {e}")
                continue

        # Process images
        for img in images:
            url = img.get("url", "")
            filename = img.get("filename", "unknown")

            if not url:
                continue

            try:
                if has_gcs and (url.startswith("gs://") or url.startswith("http")):
                    ext = Path(url).suffix.lower()
                    with download_to_temp(url, suffix=ext) as temp_path:
                        with open(temp_path, 'rb') as f:
                            img_bytes = f.read()

                        media_type = f"image/{ext[1:]}" if ext else "image/png"
                        part = types.Part.from_bytes(data=img_bytes, mime_type=media_type)
                        content_parts.append(part)
                        processed_files.append({"filename": filename, "type": "image"})
                        print(f"[Enterprise Analyzer] Loaded image: {filename}")

            except Exception as e:
                print(f"[Enterprise Analyzer] Error processing image {filename}: {e}")
                continue

        if not content_parts:
            print("[Enterprise Analyzer] No content could be processed")
            return {
                "status": "processing_failed",
                "message": "Could not process any uploaded resources",
                "timestamp": datetime.now().isoformat()
            }

        # Add analysis instruction
        content_parts.append(ENTERPRISE_ANALYSIS_PROMPT)

        # Call Gemini 3 Pro for analysis
        print(f"[Enterprise Analyzer] Calling Gemini 3 Pro with {len(content_parts)} content parts...")

        response = client.models.generate_content(
            model="gemini-3-pro-preview",
            contents=content_parts,
            config=types.GenerateContentConfig(
                temperature=1.0,
                max_output_tokens=5000,
                system_instruction=ENTERPRISE_SYSTEM_INSTRUCTION,
                thinking_config=types.ThinkingConfig(
                    thinking_level="high",
                    include_thoughts=True
                )
            )
        )

        if not response or not response.text:
            raise ValueError("Empty response from Gemini")

        # Get natural text analysis (no JSON parsing needed)
        analysis_text = response.text.strip()

        # Build final output
        output = {
            "status": "success",
            "analysis": analysis_text,
            "metadata": {
                "analyzed_files": processed_files,
                "total_documents": len(documents),
                "total_images": len(images),
                "processed_count": len(processed_files),
                "timestamp": datetime.now().isoformat(),
                "model": "gemini-3-pro-preview",
                "analysis_length": len(analysis_text)
            }
        }

        elapsed = (datetime.now() - start_time).total_seconds()
        print(f"[Enterprise Analyzer] Analysis complete in {elapsed:.2f}s")
        print(f"[Enterprise Analyzer] Processed {len(processed_files)} files")

        # Emit completion event with full output
        if state:
            from ..base import emit_event
            emit_event(state, "enterprise_analysis_complete", {
                "status": "success",
                "analysis": analysis_text,
                "metadata": output["metadata"],
                "document_count": len(documents),
                "image_count": len(images),
                "processed_count": len(processed_files),
                "analysis_length": len(analysis_text),
                "elapsed_seconds": elapsed,
                "timestamp": datetime.now().isoformat()
            }, agent_name="enterprise_analyzer")

        return output

    except Exception as e:
        import traceback
        error_trace = traceback.format_exc()
        print(f"[Enterprise Analyzer] Error: {e}")
        print(f"[Enterprise Analyzer] Traceback:\n{error_trace}")

        # Emit error event
        if state:
            from ..base import emit_event
            emit_event(state, "enterprise_analysis_error", {
                "error": str(e),
                "timestamp": datetime.now().isoformat()
            }, agent_name="enterprise_analyzer")

        return {
            "status": "error",
            "error": str(e),
            "traceback": error_trace,
            "timestamp": datetime.now().isoformat()
        }
