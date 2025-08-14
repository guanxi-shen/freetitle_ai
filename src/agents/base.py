"""
Agent implementations for AI Video Generation Studio workflow
"""

import json
import os
import time
import logging
from datetime import datetime
from typing import Dict, List, Any
from ..core.state import RAGState
from ..core.config import ROLLING_WINDOW_SIZE
from ..prompts import AGENT_REGISTRY

logger = logging.getLogger(__name__)

# Utility functions used by all agents
def emit_event(state: RAGState, event_type: str, data: Any = None, agent_name: str = None) -> RAGState:
    """
    Emit a streaming event using LangGraph's StreamWriter when available.
    Falls back to state accumulation for backward compatibility.

    Args:
        state: Current workflow state
        event_type: Type of event to emit
        data: Event data to include
        agent_name: Optional explicit agent name (takes precedence over state)

    Returns:
        Updated state (for compatibility)
    """
    try:
        # Get StreamWriter for real-time streaming
        from langgraph.config import get_stream_writer
        writer = get_stream_writer()

        # Emit event directly via StreamWriter
        # Prefer explicit agent_name parameter, fallback to state, then "unknown"
        writer({
            "type": event_type,
            "timestamp": datetime.now().isoformat(),
            "agent": agent_name or state.get("current_agent", "unknown"),
            "data": data
        })
    except Exception as e:
        # Log streaming failures for debugging
        # Only expected when called outside LangGraph context (e.g., from thread)
        print(f"[emit_event ERROR] Streaming disabled or error for event_type={event_type}, agent={agent_name or state.get('current_agent', 'unknown')}: {type(e).__name__}: {str(e)}")
        pass

    return state

def format_time(seconds: float) -> str:
    """Format time in seconds to minutes and seconds"""
    minutes = int(seconds // 60)
    remaining_seconds = seconds % 60
    if minutes > 0:
        return f"{minutes}m {remaining_seconds:.1f}s"
    else:
        return f"{remaining_seconds:.1f}s"

# Note: capture_prompt and capture_context_content functions removed (were no-ops)

def clean_json_response(response_text: str) -> str:
    """
    Clean JSON response from markdown code blocks and extra trailing braces.
    LLMs may wrap JSON in ```json...``` blocks which need to be stripped.
    Some models occasionally add extra closing braces at the end.

    Args:
        response_text: Raw text that might contain markdown-wrapped JSON

    Returns:
        Clean JSON string ready for parsing
    """
    if response_text.startswith("```json"):
        # Remove opening ```json and closing ```
        response_text = response_text[7:]  # Remove ```json
        if response_text.endswith("```"):
            response_text = response_text[:-3]  # Remove closing ```
        response_text = response_text.strip()
    elif response_text.startswith("```"):
        # Handle case where it's just ```
        response_text = response_text[3:]
        if response_text.endswith("```"):
            response_text = response_text[:-3]
        response_text = response_text.strip()

    # Remove extra trailing closing braces
    # Only remove if there's a clear imbalance (more } than {)
    # Limit to 3 removals to be safe
    max_removals = 3
    removals = 0
    while removals < max_removals:
        open_count = response_text.count('{')
        close_count = response_text.count('}')

        # If balanced or has more opens, stop
        if close_count <= open_count:
            break

        # If ends with }, remove it
        if response_text.rstrip().endswith('}'):
            response_text = response_text.rstrip()[:-1]  # Remove last }
            removals += 1
        else:
            break  # Extra } is not at the end, don't touch it

    return response_text.strip()


def load_enterprise_documents(state: RAGState, agent_name: str = "unknown") -> List[Any]:
    """
    Load enterprise documents from state in universal format

    Processes PDFs, PowerPoint, and Word documents from enterprise_resources.
    PDFs are loaded as binary data in universal format, PPT/Word are text-extracted.

    Args:
        state: RAGState containing enterprise_resources
        agent_name: Name of calling agent (for logging)

    Returns:
        List of dicts with format {"data": bytes, "media_type": str} and text strings
    """
    enterprise_resources = state.get("enterprise_resources", {})
    documents = enterprise_resources.get("documents", [])
    document_parts = []

    if not documents:
        return document_parts

    print(f"[{agent_name}] Loading {len(documents)} enterprise documents...")

    # Check GCS download utility availability
    try:
        from ...storage.gcs_utils import download_to_temp
        has_gcs = True
    except ImportError:
        has_gcs = False
        print(f"[{agent_name}] Warning: GCS utils not available for document loading")
        return document_parts

    # Process each document
    for doc in documents:
        url = doc.get("url", "")
        filename = doc.get("filename", "unknown")
        extension = doc.get("extension", "")

        if not url:
            continue

        try:
            if has_gcs and (url.startswith("gs://") or url.startswith("http")):
                with download_to_temp(url, suffix=extension) as temp_path:
                    with open(temp_path, 'rb') as f:
                        doc_bytes = f.read()

                    # PDF: Load as universal format dict
                    if extension == ".pdf":
                        document_parts.append({
                            "data": doc_bytes,
                            "media_type": "application/pdf"
                        })
                        print(f"[{agent_name}] Loaded PDF: {filename}")

                    # PowerPoint: Extract text
                    elif extension in [".pptx", ".ppt"]:
                        try:
                            from pptx import Presentation
                            import tempfile

                            with tempfile.NamedTemporaryFile(suffix=extension, delete=False) as tmp:
                                tmp.write(doc_bytes)
                                tmp_path = tmp.name

                            prs = Presentation(tmp_path)
                            text_content = []
                            for slide in prs.slides:
                                for shape in slide.shapes:
                                    if hasattr(shape, "text"):
                                        text_content.append(shape.text)

                            import os
                            os.unlink(tmp_path)

                            document_parts.append("\n\n".join(text_content))
                            print(f"[{agent_name}] Extracted text from PPT: {filename}")

                        except ImportError:
                            print(f"[{agent_name}] python-pptx not available, skipping {filename}")

                    # Word: Extract text
                    elif extension in [".docx", ".doc"]:
                        try:
                            from docx import Document
                            import tempfile

                            with tempfile.NamedTemporaryFile(suffix=extension, delete=False) as tmp:
                                tmp.write(doc_bytes)
                                tmp_path = tmp.name

                            doc_obj = Document(tmp_path)
                            text_content = [para.text for para in doc_obj.paragraphs]

                            import os
                            os.unlink(tmp_path)

                            document_parts.append("\n\n".join(text_content))
                            print(f"[{agent_name}] Extracted text from Word: {filename}")

                        except ImportError:
                            print(f"[{agent_name}] python-docx not available, skipping {filename}")

        except Exception as e:
            print(f"[{agent_name}] Error processing document {filename}: {e}")
            continue

    if document_parts:
        print(f"[{agent_name}] Successfully loaded {len(document_parts)} document parts")

    return document_parts


# Metadata extraction utility functions
def get_character_images(state: RAGState) -> List[Dict[str, Any]]:
    """Extract character images from metadata (single source of truth)

    Returns list of dicts with url, character, type, filename
    """
    images = []
    char_metadata = state.get("character_image_metadata", {})

    for filename, metadata in char_metadata.items():
        path = metadata.get("path")
        if metadata.get("success") and path:
            images.append({
                "url": path,
                "character": metadata.get("character_name", "Unknown"),
                "type": metadata.get("image_type", "unknown"),
                "filename": filename
            })

    return images


def get_storyboard_frames(state: RAGState) -> List[Dict[str, Any]]:
    """Extract storyboard frames from metadata (single source of truth)

    Returns list of dicts with url, scene, shot, frame, filename
    """
    frames = []
    frame_metadata = state.get("storyboard_frame_metadata", {})

    for filename, metadata in frame_metadata.items():
        if metadata.get("success") and metadata.get("path"):
            frames.append({
                "url": metadata["path"],
                "scene": metadata.get("scene"),
                "shot": metadata.get("shot"),
                "frame": metadata.get("frame"),
                "filename": filename
            })

    return frames


def get_supplementary_content(state: RAGState) -> List[Dict[str, Any]]:
    """Extract supplementary content from metadata (single source of truth)

    Returns list of dicts with url, content_type, content_name
    """
    content = []
    supp_metadata = state.get("supplementary_content_metadata", {})

    for key, metadata in supp_metadata.items():
        if metadata.get("success") and metadata.get("path"):
            content.append({
                "url": metadata["path"],
                "content_type": metadata.get("content_type", "unknown"),
                "content_name": metadata.get("content_name", key)
            })

    return content


def get_all_asset_urls(state: RAGState) -> Dict[str, List[str]]:
    """Get all asset URLs organized by type

    Returns dict with keys: characters, storyboards, supplementary, references, videos, audio
    """
    assets = {
        "characters": [],
        "storyboards": [],
        "supplementary": [],
        "references": [],
        "videos": [],
        "audio": []
    }

    # Characters from metadata
    for img in get_character_images(state):
        assets["characters"].append(img["url"])

    # Storyboards from generated_storyboards
    for frame in get_storyboard_frames(state):
        assets["storyboards"].append(frame["url"])

    # Supplementary from metadata
    for item in get_supplementary_content(state):
        assets["supplementary"].append(item["url"])

    # User references
    assets["references"] = state.get("reference_images", [])

    # Videos from asset_urls (source of truth)
    generated_videos = state.get("asset_urls", {}).get("generated_videos", [])
    for video in generated_videos:
        if video.get("url"):
            assets["videos"].append(video["url"])

    # Audio from asset_urls (following video pattern)
    generated_audio = state.get("asset_urls", {}).get("generated_audio", [])
    for track in generated_audio:
        if track.get("url"):
            assets["audio"].append(track["url"])

    return assets


# Define full access agents list once at module level
# Full access = agents that need shot-level CONTENT (dialogue, actions, transitions) or PRODUCTION details (compositions, progressions)
# Filtered access = agents that only need metadata (counts, summaries, character lists, production notes)
# Note: video_editor_agent removed - only needs production_notes and shot continuity notes (extracted separately)
# Note: orchestrator and answer_parser removed - only need high-level summary for routing/synthesis
FULL_ACCESS_AGENTS = ["script_agent", "storyboard_agent", "supplementary_agent", "video_agent"]

# Agents that should exclude extra metadata (character, storyboard, supplementary, video generation tasks)
# These agents only need core content (scripts, video/audio URLs) without supporting metadata
EXCLUDE_EXTRA_METADATA_AGENTS = ["video_editor_agent"]

def get_filtered_script_data(script_data: dict, agent_name: str = None) -> dict:
    """Filter script data based on agent permissions
    
    Args:
        script_data: Full script data dictionary
        agent_name: Name of the agent requesting the data
        
    Returns:
        Filtered script data - full for script_agent/storyboard_agent, summary for others
    """
    if not script_data:
        return {}
    
    if agent_name in FULL_ACCESS_AGENTS:
        # Return full script data
        return script_data
    
    # For all other agents, return high-level summary only
    script_details = script_data.get("script_details", {})
    scenes = script_details.get("scenes", [])

    # Extract character list and count dialogue/shots
    all_characters = set()
    total_shots = 0
    total_dialogue_lines = 0
    scene_summaries = []

    for scene in scenes:
        characters = scene.get("characters", [])
        all_characters.update(characters)
        shots = scene.get("shots", [])
        total_shots += len(shots)

        # Count dialogue in this scene
        scene_dialogue_count = 0
        for shot in shots:
            dialogue = shot.get("dialogue", [])
            total_dialogue_lines += len(dialogue)
            scene_dialogue_count += len(dialogue)

        # Build scene summary (no shot details)
        scene_summaries.append({
            "scene_number": scene.get("scene_number", len(scene_summaries) + 1),
            "scene_summary": scene.get("scene_summary", ""),
            "setting": scene.get("setting", "Unknown"),
            "duration": scene.get("duration", "Unknown"),
            "characters": characters,
            "shot_count": len(shots),
            "dialogue_line_count": scene_dialogue_count
        })

    # Build filtered data - high-level info only
    filtered_data = {
        "script_details": {
            "title": script_details.get("title", "Untitled"),
            "duration": script_details.get("duration", "Unknown"),
            "video_summary": script_details.get("video_summary", ""),
            "creative_vision": script_details.get("creative_vision", ""),
            "aspect_ratio": script_details.get("aspect_ratio", "horizontal"),
            "total_scenes": len(scenes),
            "total_shots": total_shots,
            "total_dialogue_lines": total_dialogue_lines,
            "characters_in_script": sorted(list(all_characters)),
            "scenes": scene_summaries  # Scene-level summaries only, no shot details
        },
        "characters": [{"name": c.get("name", "Unknown")} for c in script_data.get("characters", [])],  # Names only
        "production_notes": script_data.get("production_notes", {}),
        "audio_design": script_data.get("audio_design", {})
    }

    return filtered_data

def build_full_context(
    state: RAGState,
    agent_name: str,
    user_query: str = None,
    instruction: str = None,
    agent_specific_context: str = None,
    **context_flags
) -> str:
    """
    Build complete dynamic context block for 100% static templates

    Structure (in order):
    1. Base agent context (conversation, assets, state)
    2. Agent-specific dynamic sections (tool selection, planning guidance)
    3. Current request (user query, instruction)
    4. Instruction reminder (re-ground agent to follow guidelines)

    Args:
        state: Current RAG state
        agent_name: Name of requesting agent
        user_query: User's current query
        instruction: Task-level instruction
        agent_specific_context: Optional agent-specific sections (tool selection, planning, etc.)
        **context_flags: Passed through to get_agent_context()

    Returns:
        Complete formatted context string
    """
    context_parts = []

    # 1. Base agent context (conversation, assets, state)
    base_context = get_agent_context(state, agent_name=agent_name, **context_flags)
    context_parts.append(base_context)

    # 2. Agent-specific dynamic sections
    if agent_specific_context:
        context_parts.append(agent_specific_context)

    # 3. Current request section
    request_parts = []
    if user_query:
        request_parts.append(f"User Query: {user_query}")
    if instruction:
        request_parts.append(f"Task Instruction: {instruction}")

    if request_parts:
        current_request = "\n\n### Current Request:\n" + "\n".join(request_parts)
        context_parts.append(current_request)

    # 4. Reminder to review instructions
    instruction_reminder = "\n\n### Important:\nReview the task instructions and guidelines provided at the beginning of this prompt before proceeding. Ensure your response follows all specified rules, formats, and requirements."
    context_parts.append(instruction_reminder)

    return "\n\n".join(context_parts)

def get_agent_context(
    state: RAGState,
    agent_name: str = None,
    context_summary: bool = False,
    window_memory: bool = False,
    conversation_history: int = 0,  # 0=none, N=last N turns, -1=all
    include_generated_content: bool = False,  # Include previously generated scripts/characters
    include_reference_images: bool = False,  # Include user-uploaded reference images
    include_image_annotations: bool = False,  # Include AI-generated annotations for reference images
    include_generated_assets: bool = False,  # Include list of all generated files (characters, supplementary)
    include_tool_selections: bool = False,  # Include user tool selections
    include_expert_info: bool = False,  # Include active expert agents
    include_enterprise_analysis: bool = True,  # Include enterprise document analysis (default True)
    include_metadata: bool = True  # Include detailed asset metadata (character images, storyboards, videos, audio)
) -> str:
    """Smart context builder for video production agents

    Args:
        state: Current RAG state containing all creative conversation and resource data
        agent_name: Name of the agent requesting context (for permission filtering)
        context_summary: Include high-level conversation context summary
        window_memory: Include recent conversation window (last N turns based on ROLLING_WINDOW_SIZE)
        conversation_history: Extended conversation access (0=none, N=last N turns, -1=all)
        include_generated_content: Include previously generated scripts/characters
        include_reference_images: Include metadata about user-uploaded reference images
        include_image_annotations: Include AI-generated annotations for reference images
        include_generated_assets: Include list of all generated files
        include_tool_selections: Include user tool selections for agents
        include_expert_info: Include active expert agents
        include_enterprise_analysis: Include enterprise document analysis (default True)

    Returns:
        Formatted context string for agent consumption

    Note:
        - conversation_history excludes window_memory turns to avoid duplication
        - reference_images provides text metadata only - agents don't see actual images
        - enterprise_analysis is included by default for all agents
    """
    content_parts = ["Context information: "]
    
    # Always include basic information
    user_query = state.get("user_query", "")
    turn_number = state.get("turn_number", 1)
    retry_count = state.get("retry_count", 0)
    
    # Basic context (always provided)
    basic_info = f"### Workflow Context:\n"
    # User query removed - appears in "Current Request" section below (build_full_context)
    basic_info += f"Turn Number: {turn_number}\n"
    if retry_count > 0:
        basic_info += f"Retry Attempt: {retry_count + 1}\n"

    # Add workflow information if available (system agents only)
    # Shows previous orchestrator run info on retries within same turn
    if agent_name in ["orchestrator", "answer_parser"]:
        selected_agents = state.get("selected_agents", [])
        execution_plan = state.get("execution_plan", "")
        if selected_agents or execution_plan:
            # retry_count reflects current attempt, so previous run was attempt {retry_count}
            basic_info += f"\n### Previous Run (Turn {turn_number}, Attempt {retry_count}):\n"
            if selected_agents:
                # Handle nested lists for parallel execution
                formatted_agents = []
                for item in selected_agents:
                    if isinstance(item, list):
                        formatted_agents.append(f"[{', '.join(item)}]")
                    else:
                        formatted_agents.append(item)
                basic_info += f"Selected Agents: {', '.join(formatted_agents)}\n"
            if execution_plan:
                basic_info += f"Execution Plan: {execution_plan}\n"
        # Old code without "Previous Run" header:
        # if selected_agents:
        #     formatted_agents = []
        #     for item in selected_agents:
        #         if isinstance(item, list):
        #             formatted_agents.append(f"[{', '.join(item)}]")
        #         else:
        #             formatted_agents.append(item)
        #     basic_info += f"Selected Agents: {', '.join(formatted_agents)}\n"
        # if execution_plan:
        #     basic_info += f"Execution Plan: {execution_plan}\n"
    
    content_parts.append(basic_info)

    # System capabilities (concise format)
    capabilities_info = []

    # Tool selections
    if include_tool_selections:
        tool_selections = state.get("user_preferences", {}).get("tool_selections", {})
        if tool_selections:
            tool_parts = [f"{agent}=[{', '.join(tools)}]" for agent, tools in tool_selections.items()]
            capabilities_info.append(f"Tool Selections: {', '.join(tool_parts)}")

        # Workflow mode (only show if do-it-all)
        workflow_mode = state.get("user_preferences", {}).get("workflow_mode", "step-by-step")
        if workflow_mode == "do-it-all":
            capabilities_info.append(
                "Workflow Mode: do-it-all\n"
                "- User prefers automatic creative execution with minimal conversation\n"
                "- You have full creative control to make production decisions\n"
                "- Use parallel execution for independent agents (e.g., \"parallel:character_agent,supplementary_agent\" and \"parallel:video_agent,audio_agent\") for faster results\n"
                "- Use multi-agent chaining, proceed through production stages autonomously using best judgment\n"
                "- Still separate pre-production (Script, Character, Supplement, Storyboard) and production (Video, Audio, Editing) with a confirmation response in between the two stages, don't chain the two stages together unless user specifically asks\n"
                "- Override default follow-up policy, no greeting/follow-up/clarification questions needed unless critically necessary for completion"
            )

    # Active experts
    if include_expert_info:
        expert_selections = state.get("expert_selections", [])
        if expert_selections:
            capabilities_info.append(f"Active Experts: {', '.join(expert_selections)}")

    if capabilities_info:
        content_parts.append("### System Capabilities:\n" + "\n".join(capabilities_info))

    # Context summary (if requested and available)
    if context_summary and state.get("context_summary"):
        content_parts.append(f"### Conversation Context Summary:\n{state['context_summary']}")
    elif context_summary:
        content_parts.append("### Conversation Context Summary:\nNo previous conversation summary available (first conversation)")
    
    # Window memory (recent conversations)
    if window_memory and state.get("window_memory"):
        window_text = "### Recent Conversation History:\n"
        for conv in state["window_memory"]:
            window_text += f"Turn {conv.get('turn_number', 0)}:\nUSER: {conv.get('query', '')}\nSYSTEM: {conv.get('answer', '')}\n"
        content_parts.append(window_text)

        # NEW: Include expert window memory (windowed by global turn)
        expert_window = state.get("expert_window_memory", {})

        if expert_window:
            # Build chronological view of all expert interactions in window
            all_expert_convs = []
            for expert_id, conversations in expert_window.items():
                for conv in conversations:
                    conv_with_id = dict(conv)
                    conv_with_id['expert_id'] = expert_id
                    all_expert_convs.append(conv_with_id)

            # Sort by global turn to maintain chronological order
            all_expert_convs.sort(key=lambda x: x.get('global_turn', 0))

            if all_expert_convs:
                expert_conv_text = "\n### Recent Expert Interactions (chronological):\n"
                for conv in all_expert_convs:
                    source_tag = f" [{conv.get('source', '')}]" if conv.get('source') else ""
                    expert_conv_text += f"Global Turn {conv.get('global_turn', 0)}{source_tag}:\n"
                    expert_conv_text += f"USER: {conv.get('user_query', '')}\n"
                    expert_conv_text += f"[{conv.get('expert_id', 'unknown')}]: {conv.get('expert_response', '')}\n\n"

                content_parts.append(expert_conv_text)
    elif window_memory:
        content_parts.append("### Recent Conversation History:\nNo recent conversation history available (first conversation)")

    # Extended conversation history - exclude window memory conversations
    if conversation_history != 0:
        full_history = state.get("conversation_history", [])
        window_size = len(state.get("window_memory", []))
        
        if window_memory and window_size > 0:
            # Skip the window memory conversations
            history_to_show = full_history[:-window_size] if len(full_history) > window_size else []
        else:
            # No window memory requested, show from full history
            history_to_show = full_history
        
        if conversation_history != -1:
            history_to_show = history_to_show[-conversation_history:] if history_to_show else []
        
        if history_to_show:
            hist_text = "### Extended Conversation History:\n"
            for conv in history_to_show:
                hist_text += f"Turn {conv.get('turn_number', 0)}:\nUSER: {conv.get('query', '')}\nSYSTEM: {conv.get('answer', '')}...\n"
            content_parts.append(hist_text)
        else:
            content_parts.append("### Extended Conversation History:\nNo extended conversation history available")

    # Include previously generated content if requested
    if include_generated_content:
        script_data = state.get("generated_scripts", {})

        if script_data:
            scripts_text = "### Current Generated Script:\n"

            if agent_name in FULL_ACCESS_AGENTS:
                # Full access agents get complete unfiltered script data
                scripts_text += "Full script with all scene and shot details:\n\n"
                scripts_text += json.dumps(script_data, indent=2, ensure_ascii=False)
                scripts_text += "\n"
            else:
                # Limited access agents get filtered summary (shot details omitted)
                filtered_script = get_filtered_script_data(script_data, agent_name)
                script_details = filtered_script.get("script_details", {})

                # Add context summary before JSON
                total_scenes = script_details.get('total_scenes', 0)
                total_shots = script_details.get('total_shots', 0)
                total_dialogue = script_details.get('total_dialogue_lines', 0)

                scripts_text += f"FILTERED SCRIPT (Intentional - complete script exists, you receive high-level summary only for efficiency):\n"
                scripts_text += f"Shot-level details omitted.\n"
                scripts_text += f"Summary: {total_scenes} scenes with {total_shots} total shots and {total_dialogue} dialogue lines.\n\n"
                scripts_text += json.dumps(filtered_script, indent=2, ensure_ascii=False)
                scripts_text += "\n"

            content_parts.append(scripts_text)

            # Add character profiles from script (only for character_agent)
            characters_profiles = script_data.get("characters", [])
            if characters_profiles and agent_name == "character_agent":
                chars_text = "### Character Profiles:\n"
                for char in characters_profiles:
                    char_name = char.get('name', 'Unknown')
                    chars_text += f"- {char_name}: {char.get('attributes', 'No description')} ({char.get('role', 'unknown role')})\n"
                content_parts.append(chars_text)

        # Show character image metadata if available for name coordination
        if include_metadata and agent_name not in EXCLUDE_EXTRA_METADATA_AGENTS:
            character_metadata = state.get("character_image_metadata", {})
            if character_metadata:
                # Fields to exclude from character metadata
                character_exclude_fields = ["params", "thinking"]  # Duplicates prompt information, internal reasoning

                filtered_char_metadata = {}
                for filename, meta in character_metadata.items():
                    if isinstance(meta, dict):
                        filtered_char_metadata[filename] = {
                            k: v for k, v in meta.items()
                            if k not in character_exclude_fields
                        }
                    else:
                        filtered_char_metadata[filename] = meta

                char_meta_text = "\n### Character Image Metadata (use exact names for consistency):\n"
                char_meta_text += json.dumps(filtered_char_metadata, indent=2, default=str)
                content_parts.append(char_meta_text)

        # Show storyboard frame metadata if available
        if include_metadata and agent_name not in EXCLUDE_EXTRA_METADATA_AGENTS:
            storyboard_metadata = state.get("storyboard_frame_metadata", {})
            if storyboard_metadata:
                # Fields to exclude from storyboard metadata
                storyboard_exclude_fields = ["thinking", "operations", "orchestrator_config", "params"]

                filtered_metadata = {}
                for filename, meta in storyboard_metadata.items():
                    if isinstance(meta, dict):
                        filtered_metadata[filename] = {
                            k: v for k, v in meta.items()
                            if k not in storyboard_exclude_fields
                        }
                    else:
                        filtered_metadata[filename] = meta

                story_meta_text = "\n### Storyboard Frame Metadata:\n"
                story_meta_text += json.dumps(filtered_metadata, indent=2, default=str)
                content_parts.append(story_meta_text)

        # # OLD: Show supplementary content metadata separately (COMMENTED OUT - now merged below)
        # if include_metadata:
        #     supplementary_content_meta = state.get("supplementary_content_metadata", {})
        #     if supplementary_content_meta:
        #         # Fields to exclude from supplementary metadata
        #         supplementary_exclude_fields = ["params", "thinking"]
        #
        #         filtered_supp_metadata = {}
        #         for filename, meta in supplementary_content_meta.items():
        #             if isinstance(meta, dict):
        #                 filtered_supp_metadata[filename] = {
        #                     k: v for k, v in meta.items()
        #                     if k not in supplementary_exclude_fields
        #                 }
        #             else:
        #                 filtered_supp_metadata[filename] = meta
        #
        #         supp_meta_text = "\n### Supplementary Content Metadata:\n"
        #         supp_meta_text += json.dumps(filtered_supp_metadata, indent=2, default=str)
        #         content_parts.append(supp_meta_text)

        # Merged supplementary content display (combines technical metadata + semantic data)
        if (include_metadata or include_generated_content) and agent_name not in EXCLUDE_EXTRA_METADATA_AGENTS:
            supp_metadata_raw = state.get("supplementary_content_metadata", {})
            supp_semantic = state.get("generated_supplementary", {})

            # Filter supplementary metadata to exclude internal fields
            supplementary_exclude_fields = ["params", "thinking"]
            supp_metadata = {}
            for filename, meta in supp_metadata_raw.items():
                if isinstance(meta, dict):
                    supp_metadata[filename] = {
                        k: v for k, v in meta.items()
                        if k not in supplementary_exclude_fields
                    }
                else:
                    supp_metadata[filename] = meta

            if supp_metadata or supp_semantic:
                supp_text = "\n### Supplementary Content:\n"

                # Organize by category from semantic data
                categories = {}
                for item_key, sem_data in supp_semantic.items():
                    category = sem_data.get("category", "uncategorized")
                    if category not in categories:
                        categories[category] = []
                    categories[category].append((item_key, sem_data))

                # Display by category with merged technical + semantic data
                for category, items in sorted(categories.items()):
                    supp_text += f"**{category.upper()}:**\n\n"

                    for item_key, sem_data in items:
                        # Find matching technical metadata by content_name
                        tech_data = None
                        for filename, meta in supp_metadata.items():
                            if meta.get('content_name') == item_key or filename.replace(".png", "") == item_key:
                                tech_data = meta
                                break

                        # Display merged entry
                        title = sem_data.get('title', item_key)
                        supp_text += f"  • {title}\n"

                        # Semantic fields
                        if sem_data.get('description'):
                            supp_text += f"    Description: {sem_data['description']}\n"
                        if sem_data.get('usage_notes'):
                            supp_text += f"    Usage: {sem_data['usage_notes']}\n"
                        if sem_data.get('related_to'):
                            supp_text += f"    Related to: {', '.join(sem_data['related_to'])}\n"

                        # Technical fields (if available)
                        path = sem_data.get('image_path') or (tech_data.get('path') if tech_data else None)
                        if path:
                            supp_text += f"    Path: {path}\n"

                        if tech_data:
                            if tech_data.get('success') is False:
                                supp_text += f"    Status: Failed\n"
                            if tech_data.get('timestamp'):
                                supp_text += f"    Generated: {tech_data['timestamp']}\n"
                            # # Tool info - COMMENTED OUT: params field filtered out for consistency
                            # if tech_data.get('params', {}).get('image_tool'):
                            #     supp_text += f"    Tool: {tech_data['params']['image_tool']}\n"

                        supp_text += "\n"

                content_parts.append(supp_text)

        # Show audio generation metadata if available
        if include_metadata:
            audio_metadata = state.get("audio_generation_metadata", {})
            if audio_metadata:
                # Fields to exclude from audio metadata
                audio_exclude_fields = ["params", "thinking"]

                filtered_audio_metadata = {}
                for filename, meta in audio_metadata.items():
                    if isinstance(meta, dict):
                        filtered_audio_metadata[filename] = {
                            k: v for k, v in meta.items()
                            if k not in audio_exclude_fields
                        }
                    else:
                        filtered_audio_metadata[filename] = meta

                audio_meta_text = "\n### Audio Generation Metadata:\n"
                audio_meta_text += json.dumps(filtered_audio_metadata, indent=2, default=str)
                content_parts.append(audio_meta_text)
        
        # Include video production data with agent-appropriate detail level
        video_prompts = state.get("generated_video_prompts", {})
        if video_prompts:
            # video_agent needs prompts for get_available_shots() tool
            if agent_name == "video_agent":
                video_prompts_text = "\n### Generated Video Prompts:\n"
                video_prompts_text += f"Video prompts with motion descriptions and frame references:\n"
                video_prompts_text += json.dumps(video_prompts, indent=2, default=str)
                content_parts.append(video_prompts_text)
            # video_editor_agent excluded - builds video list from asset_urls, doesn't need generation prompts

        if include_metadata and agent_name == "answer_parser":
            video_tasks = state.get("video_generation_tasks", [])
            if video_tasks:
                # Fields to exclude from video tasks (top-level and nested)
                video_task_exclude_fields = ["error_context"]
                video_task_metadata_exclude_fields = [
                    "thinking",
                    "reasoning",
                    "generation_prompt",  # Duplicates params.generation_prompt
                    "parameters"  # Duplicates entire params object
                ]

                filtered_tasks = []
                for task in video_tasks:
                    if isinstance(task, dict):
                        # Filter top-level fields
                        filtered_task = {k: v for k, v in task.items() if k not in video_task_exclude_fields}

                        # Filter nested metadata fields
                        if "metadata" in filtered_task and isinstance(filtered_task["metadata"], dict):
                            filtered_task["metadata"] = {
                                k: v for k, v in filtered_task["metadata"].items()
                                if k not in video_task_metadata_exclude_fields
                            }

                        filtered_tasks.append(filtered_task)
                    else:
                        filtered_tasks.append(task)

                video_tasks_text = "\n### Video Generation Tasks:\n"
                video_tasks_text += json.dumps(filtered_tasks, indent=2, default=str)
                content_parts.append(video_tasks_text)

        edited_videos = state.get("edited_videos", [])
        if edited_videos:
            # Fields to exclude from edited videos
            edited_video_exclude_fields = [
                "operations_log",
                "llm_thinking",
                "llm_full_response",
                "video_analysis_data",
                "audio_analysis_data"
            ]

            filtered_videos = []
            for video in edited_videos:
                if isinstance(video, dict):
                    filtered_videos.append({
                        k: v for k, v in video.items()
                        if k not in edited_video_exclude_fields
                    })
                else:
                    filtered_videos.append(video)

            edited_videos_text = "\n### Edited Videos:\n"
            edited_videos_text += json.dumps(filtered_videos, indent=2, default=str)
            content_parts.append(edited_videos_text)
    
    # Include user-uploaded reference images with annotations
    if include_reference_images:
        user_reference_images = state.get("reference_images", [])
        image_annotations = state.get("image_annotations", {})

        if user_reference_images and image_annotations:
            reference_text = "### Reference Images:\n\n"

            for url in user_reference_images:
                # Extract filename key (annotations are keyed by filename, not full URL)
                filename = url.split("/")[-1].split("?")[0] if "/" in url else url

                # Match by filename key
                annotation_data = image_annotations.get(filename, {})
                description = annotation_data.get("description", "No annotation available")

                reference_text += f"**{filename}**\n"
                reference_text += f"Path: {url}\n"
                reference_text += f"AI annotation: {description}\n\n"

            content_parts.append(reference_text)

    # # OLD: Verbose reference image section (DEPRECATED)
    # if include_reference_images:
    #     user_reference_images = state.get("reference_images", [])
    #     if user_reference_images:
    #         image_annotations = state.get("image_annotations", {})
    #         reference_text = "### User-Uploaded Reference Images:\n"
    #         reference_text += f"The user uploaded {len(user_reference_images)} reference image(s).\n"
    #         if image_annotations:
    #             reference_text += "Each image includes an AI-generated annotation describing its visual content.\n\n"
    #         else:
    #             reference_text += "\n"
    #         for idx, url in enumerate(user_reference_images, 1):
    #             filename = url.split("/")[-1].split("?")[0] if "/" in url else url
    #             reference_text += f"{idx}. {filename}\n"
    #             reference_text += f"   Path: {url}\n"
    #             annotation_data = image_annotations.get(url, {})
    #             if annotation_data and annotation_data.get("description"):
    #                 description = annotation_data["description"]
    #                 reference_text += f"   Content: {description}\n"
    #             elif image_annotations:
    #                 reference_text += f"   Content: (annotation pending)\n"
    #             reference_text += "\n"
    #         content_parts.append(reference_text)
    
    
    # Include all generated assets for supplementary agent
    if include_generated_content and agent_name == "supplementary_agent":
        # # COMMENTED OUT: Duplicate character images simple list (already shown in metadata below)
        # character_images = list_all_character_images(state)
        # if character_images:
        #     char_text = "### Generated Character Images:\n"
        #     char_text += f"Character design reference images available for use:\n"
        #     for url in character_images:
        #         filename = url.split("/")[-1] if "/" in url else url
        #         char_text += f"  - {filename}: {url}\n"
        #     content_parts.append(char_text)

        # # COMMENTED OUT: Duplicate storyboard images list (already shown in metadata above)
        # storyboard_frames = get_storyboard_frames(state)
        # storyboard_images = [frame["url"] for frame in storyboard_frames]
        # if storyboard_images:
        #     story_text = "### Generated Storyboard Images:\n"
        #     story_text += f"Storyboard frames available for reference:\n"
        #     for url in storyboard_images:
        #         filename = url.split("/")[-1] if "/" in url else url
        #         story_text += f"  - {filename}: {url}\n"
        #     content_parts.append(story_text)

        # Character images from metadata
        character_metadata = state.get("character_image_metadata", {})
        if character_metadata:
            char_imgs_text = "### Generated Character Images:\n"
            # Group by character name
            char_groups = {}
            for filename, meta in character_metadata.items():
                if meta.get("success") and meta.get("path"):
                    char_name = meta.get("character_name", "Unknown")
                    if char_name not in char_groups:
                        char_groups[char_name] = []
                    char_groups[char_name].append(meta)

            for char_name, images in sorted(char_groups.items()):
                char_imgs_text += f"\n{char_name}:\n"
                for meta in images:
                    img_type = meta.get("image_type", "unknown")
                    path = meta.get("path", "")
                    char_imgs_text += f"  - {img_type}: {path}\n"
                    # Include prompt for visual understanding
                    if meta.get("prompt"):
                        prompt_preview = meta["prompt"][:100] + "..." if len(meta["prompt"]) > 100 else meta["prompt"]
                        char_imgs_text += f"    Prompt: {prompt_preview}\n"
            content_parts.append(char_imgs_text)

    # # OLD: Generated supplementary content (semantic only) - COMMENTED OUT - now merged with metadata above
    # if include_generated_content:
    #     supplementary_metadata = state.get("generated_supplementary", {})
    #     if supplementary_metadata:
    #         supp_text = "### Generated Supplementary Content:\n"
    #         supp_text += "Supplementary creative materials:\n\n"
    #
    #         # Organize by category
    #         categories = {}
    #         for item_key, item_data in supplementary_metadata.items():
    #             category = item_data.get("category", "uncategorized")
    #             if category not in categories:
    #                 categories[category] = []
    #             categories[category].append((item_key, item_data))
    #
    #         # Display by category
    #         for category, items in categories.items():
    #             supp_text += f"**{category.upper()}:**\n"
    #             for item_key, item_data in items:
    #                 supp_text += f"  • {item_data.get('title', item_key)}\n"
    #                 supp_text += f"    - {item_data.get('description', 'No description')}\n"
    #                 if item_data.get('usage_notes'):
    #                     supp_text += f"    - Usage: {item_data['usage_notes']}\n"
    #                 if item_data.get('image_path'):
    #                     supp_text += f"    - Image: {item_data['image_path']}\n"
    #                 if item_data.get('related_to'):
    #                     supp_text += f"    - Related to: {', '.join(item_data['related_to'])}\n"
    #                 supp_text += "\n"
    #
    #         content_parts.append(supp_text)
    
    # OLD: Detailed Image Annotations section (DEPRECATED - now merged with reference images above)
    # if include_image_annotations:
    #     image_annotations = state.get("image_annotations", {})
    #     if image_annotations:
    #         annotations_text = "### Detailed Image Annotations:\n"
    #         annotations_text += "Full AI-generated descriptions of user-uploaded reference images.\n"
    #         annotations_text += "Use these to understand what references contain and recreate their visual content.\n\n"
    #         for url, annotation_data in image_annotations.items():
    #             description = annotation_data.get("description", "No description available")
    #             display_name = url.split("/")[-1].split("?")[0] if "/" in url else url
    #             annotations_text += f"**{display_name}**\n"
    #             annotations_text += f"{description}\n"
    #             annotations_text += "---\n"
    #         content_parts.append(annotations_text)

    # Include list of all generated assets if requested
    if include_generated_assets:
        all_generated = []
        gen_text = "### Previously Generated Assets:\n"
        gen_text += "Assets already generated in this session:\n"
        
        # Extract character images from metadata (utility function)
        character_images = get_character_images(state)

        if character_images:
            gen_text += "\nCharacter Images:\n"
            for img in character_images:
                gen_text += f"  - {img['character']} ({img['type']}): {img['url']}\n"
        
        # Extract storyboard frames from metadata (utility function)
        storyboard_frames = get_storyboard_frames(state)

        if storyboard_frames:
            gen_text += "\nStoryboard Frames:\n"
            for frame in sorted(storyboard_frames, key=lambda x: (x.get("scene", 0), x.get("shot", 0), x.get("frame", 0))):
                gen_text += f"  - Scene {frame['scene']}, Shot {frame['shot']}, Frame {frame['frame']}: {frame['url']}\n"
        
        # Extract supplementary content from metadata (utility function)
        supplementary_assets = get_supplementary_content(state)

        if supplementary_assets:
            gen_text += "\nSupplementary Content:\n"
            for asset in supplementary_assets:
                gen_text += f"  - {asset['content_name']} ({asset['content_type']}): {asset['url']}\n"

        # Include generated audio tracks
        audio_metadata = state.get("audio_generation_metadata", {})
        audio_tracks = [m for m in audio_metadata.values() if m.get("success") and m.get("path")]
        if audio_tracks:
            gen_text += "\nGenerated Audio Tracks:\n"
            for track in audio_tracks:
                track_name = track.get("name", "unknown")
                duration_sec = track.get("duration_ms", 0) / 1000
                gen_text += f"  - {track_name} ({duration_sec}s): {track.get('path')}\n"

        # Include generated videos if any
        generated_videos = state.get("asset_urls", {}).get("generated_videos", [])
        if generated_videos:
            gen_text += "\nGenerated Videos:\n"
            for video in generated_videos:
                gen_text += f"  - Scene {video.get('scene', '?')}, Shot {video.get('shot', '?')} (v{video.get('version', 1)}): {video.get('url')}\n"

        if character_images or storyboard_frames or supplementary_assets or audio_tracks or generated_videos:
            gen_text += "\nNote: These are GCS URLs. Check this list before generating new content.\n"
            content_parts.append(gen_text)
    
    # Include enterprise document analysis (default True)
    if include_enterprise_analysis:
        enterprise_output = state.get("enterprise_agent_output", {})
        if enterprise_output.get("status") == "success":
            analysis = enterprise_output.get("analysis", "")
            if analysis:
                enterprise_text = "### Enterprise Document Analysis:\n"
                enterprise_text += "Business context from uploaded documents:\n\n"
                enterprise_text += analysis
                content_parts.append(enterprise_text)
        elif enterprise_output.get("status") == "error":
            # Include error info for debugging
            enterprise_text = "### Enterprise Document Analysis:\n"
            enterprise_text += f"Note: Document analysis encountered an error: {enterprise_output.get('error', 'Unknown error')}\n"
            content_parts.append(enterprise_text)

    # Inject validation instructions for multimodal agents (automatic)
    from ..core.config import ENABLE_MULTIMODAL_VALIDATION, INCLUDE_REFERENCE_IMAGES_IN_MULTIMODAL
    if ENABLE_MULTIMODAL_VALIDATION:

        if agent_name == "character_agent":
            validation = """
### Image Quality Check

After calling generation functions, first confirm you can see the generated visual output returned in the function results. If so, think through and analyze: prompt intent, generation quality, and references alignment.

Check generated images for significant issues (max 2 regeneration attempts).

Expected: Turnarounds = 4-panel grid (front, 3/4, side, back) - CORRECT format

Check for MAJOR issues:
- Generated image satisfies your needs and generation model responds to your prompts well
- Unwanted AI glitches: Distorted anatomy, faces, hands, unnatural artifacts"""

            if INCLUDE_REFERENCE_IMAGES_IN_MULTIMODAL:
                validation += """
- Reference mismatch: Products - logos accurate (shape, text, placement), brand colors precise, features visible. Characters - identity features match"""

            validation += """

Regenerate only if significant problems. Accept minor imperfections after 2 attempts.
"""
            content_parts.append(validation)
            logger.info(f"[Validation] Injected for character_agent (references: {'enabled' if INCLUDE_REFERENCE_IMAGES_IN_MULTIMODAL else 'disabled'})")

        elif agent_name == "storyboard_shot_processor":
            validation = """
### Image Quality Check

After calling generation functions, first confirm you can see the generated visual output returned in the function results. If so, think through and analyze: prompt intent, generation quality, and references alignment.

Check generated frames for significant issues (max 2 regeneration attempts).

Check for MAJOR issues:
- Generated frame satisfies your needs and generation model responds to your prompts well
- Unwanted AI glitches: Distorted faces, hands, anatomy, unnatural artifacts, blurry regions
- Unwanted layout issues: Unintended split screens, borders, frames, grids, multi-panel layouts when not requested
- Text accuracy (if visible): Spelling, readability, valid characters
- Overly simplistic style when not intended: Flat backgrounds, minimal details, empty environments (unless script specifies simplicity)
- Dual-frame within one shot: When generating 2 frames for a shot (dual-frame mode), the second frame should not be identical to the first frame (same pose/angle/composition wastes video generation)
- Background-subject coherence: Subject perspective/scale must match background vanishing point, physical grounding (shadows, ground contact) must be coherent with environment
- Dual-frame motion: If subject moves/rotates between f1 and f2, background must adapt accordingly (parallax, angle shift) - static background with moving subject is unnatural"""

            if INCLUDE_REFERENCE_IMAGES_IN_MULTIMODAL:
                validation += """
- Reference mismatch: Characters - identity features must match. Products - logos accurate (shape, text, placement), brand colors precise"""

            validation += """

Regenerate only if significant problems. Accept minor imperfections after 2 attempts.
"""
            content_parts.append(validation)
            logger.info(f"[Validation] Injected for storyboard_shot_processor (references: {'enabled' if INCLUDE_REFERENCE_IMAGES_IN_MULTIMODAL else 'disabled'})")

        elif agent_name == "supplementary_agent":
            validation = """
### Image Quality Check

After calling generation functions, first confirm you can see the generated visual output returned in the function results. If so, think through and analyze: prompt intent, generation quality, and references alignment.

Check generated content for significant issues (max 2 regeneration attempts).

Check for MAJOR issues:
- Generated content satisfies your needs and generation model responds to your prompts well
- Unwanted AI glitches: Distorted anatomy, faces, hands, unnatural artifacts"""

            if INCLUDE_REFERENCE_IMAGES_IN_MULTIMODAL:
                validation += """
- Reference mismatch: Style, mood, composition should match references"""

            validation += """

Regenerate only if significant problems. Accept minor imperfections after 2 attempts.
"""
            content_parts.append(validation)
            logger.info(f"[Validation] Injected for supplementary_agent (references: {'enabled' if INCLUDE_REFERENCE_IMAGES_IN_MULTIMODAL else 'disabled'})")

    # If no content was added beyond basic info, provide a helpful fallback
    if len(content_parts) == 1:  # Only basic info
        content_parts.append("### Additional Context:\nThis appears to be the first conversation turn with no historical context available.")

    final_context = "\n\n".join(content_parts)

    # Token counting per category (drop-in, uses flag from config)
    from ..core.config import DEBUG_TOKEN_COUNTING
    if DEBUG_TOKEN_COUNTING:
        try:
            from ..core.token_counter import count_tokens

            # Determine model based on current agent
            current_agent = agent_name or state.get("current_agent", "unknown")

            model_name = "gemini-3-pro-preview"
            model_limit = 1_048_576

            print(f"\n{'='*80}")
            print(f"[TOKEN DEBUG] get_agent_context({agent_name or 'unknown'}) - Model: {model_name}")
            print(f"{'='*80}")
            print(f"{'Section':<40} {'Tokens':>12} {'% of Limit':>12}")
            print(f"{'-'*80}")

            # Count each section
            for i, part in enumerate(content_parts):
                try:
                    tokens = count_tokens(model_name, part)
                    section_name = part.split('\n')[0] if part else f"Section {i}"
                    if section_name.startswith('###'):
                        section_name = section_name.replace('###', '').strip()[:35]
                    else:
                        section_name = f"Section {i}"
                    print(f"{section_name:<40} {tokens:>12,} {tokens*100/model_limit:>11.2f}%")
                except Exception as e:
                    print(f"[TOKEN DEBUG] Error counting section {i}: {e}")

            # Total
            try:
                total_tokens = count_tokens(model_name, final_context)
                print(f"{'-'*80}")
                print(f"{'TOTAL':<40} {total_tokens:>12,} {total_tokens*100/model_limit:>11.2f}%")
            except Exception as e:
                print(f"[TOKEN DEBUG] Error counting total: {e}")

            print(f"{'='*80}\n")
        except Exception as e:
            print(f"[TOKEN DEBUG] Error initializing token counting: {e}")

    return final_context

def get_agent_info(agent_name: str) -> str:
    """Get formatted agent information for orchestrator"""
    if agent_name in AGENT_REGISTRY:
        agent = AGENT_REGISTRY[agent_name]
        return f"{agent_name}: {agent['description']} (Capabilities: {', '.join(agent['capabilities'])})"
    return f"{agent_name}: Unknown agent"

# Reference management functions removed - no longer needed with GCS-only approach
# References are now managed as GCS URLs, not local files

# organize_reference removed - no longer needed with GCS URLs

def list_all_references(state) -> Dict[str, List[str]]:
    """List all reference images from state
    
    Args:
        state: Current workflow state (can be RAGState or dict)
        
    Returns:
        Dictionary with 'user_references' key containing GCS URLs
    """
    references = {}
    
    # Extract user reference images from state
    user_refs = state.get("reference_images", []) if hasattr(state, "get") else []
    if user_refs:
        references["user_references"] = user_refs
    
    return references

def list_all_character_images(state) -> List[str]:
    """List all generated character images from state

    Args:
        state: Current workflow state (can be RAGState or dict)

    Returns:
        List of GCS URLs for all character images
    """
    character_images = []

    # Extract from character_image_metadata (single source of truth)
    char_metadata = state.get("character_image_metadata", {}) if hasattr(state, "get") else {}
    for filename, metadata in char_metadata.items():
        path = metadata.get("path")
        if metadata.get("success") and path:
            character_images.append(path)

    return character_images

def list_all_supplementary_images(state) -> List[str]:
    """List all generated supplementary content images from state

    Args:
        state: Current workflow state (can be RAGState or dict)

    Returns:
        List of GCS URLs for all supplementary images
    """
    supplementary_images = []

    # Extract from supplementary_content_metadata (single source of truth)
    supp_metadata = state.get("supplementary_content_metadata", {}) if hasattr(state, "get") else {}
    for key, metadata in supp_metadata.items():
        if metadata.get("success") and metadata.get("path"):
            supplementary_images.append(metadata["path"])

    return supplementary_images


# Export all functions for backward compatibility
__all__ = [
    # Utility functions
    'format_time',
    'get_agent_context',
    'get_filtered_script_data',
    'get_agent_info',
    'emit_event',
    'load_enterprise_documents',

    # Reference management functions
    'list_all_references',
    'list_all_character_images',
    'list_all_supplementary_images'
]